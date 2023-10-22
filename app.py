from flask import Flask, render_template, request, send_file
from gensim.summarization import summarize
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
import tempfile
import random

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

# @app.route('/summarize', methods=['POST'])
# def summarize_text():
#     text = request.form['text']
#     summary = summarize(text)
#     return render_template('summary.html', summary=summary, original_text=text)

@app.route('/upload', methods=['POST'])
def success():
    if request.method == 'POST':
        f = request.files['file']
        f.save(f.filename)

        # Set the path to your PDF file and the output PPTX file in a temporary directory
        pdf_path = f.filename
        output_pptx_path = tempfile.mktemp(suffix=".pptx")

        create_ppt_from_summarized_content(pdf_path, output_pptx_path)

        return send_file(output_pptx_path, as_attachment=True, download_name="output.pptx")

def extract_headings_and_content(pdf_path):
    headings_and_content = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            page_text = page.extract_text()
            if page_text:
                lines = page_text.split('\n')
                headings = []
                content = []
                for line in lines:
                    if line.isupper() or line.startswith('Chapter') or line.startswith('Section'):
                        headings.append(line)
                    else:
                        content.append(line)
                if headings:
                    headings_and_content.append((headings, content))
    return headings_and_content

def create_ppt_from_summarized_content(pdf_path, output_pptx_path):
    headings_and_content = extract_headings_and_content(pdf_path)
    presentation = Presentation()
    
    # Define the maximum number of slides to generate
    max_slides = 5  # You can adjust this as needed
    
    # Shuffle the order of content for randomness
    random.shuffle(headings_and_content)
    
    for i, (headings, content) in enumerate(headings_and_content):
        if i >= max_slides:
            break  # Stop after generating the maximum number of slides
        
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
        if headings:
            title = slide.shapes.title
            title.text = headings[0]  # Use the first heading as the slide title
        if content:
            word_count = random.randint(100, 150)
            summarized_content = summarize("\n".join(content), word_count=word_count)            
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = summarized_content
    presentation.save(output_pptx_path)

if __name__ == '__main__':
    app.run(debug=True)
